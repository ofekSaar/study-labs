import fitz  # PyMuPDF
import io
from PIL import Image
import pytesseract
from pptx import Presentation

def extract_text_from_bytes(file_bytes, filename=None):
    """
    Extracts text from PDF or PPTX bytes.
    Handles standard text-based PDFs and falls back to OCR for image-based pages.
    Extracts text from PPTX slides.
    """
    print(f"OCR Service: Processing {filename or 'bytes'}...")
    text_content = []
    
    if filename and filename.lower().endswith('.pptx'):
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text_content.append(f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text))
            return "\n\n".join(text_content)
        except Exception as e:
            print(f"PPTX Extraction Error: {e}")
            return ""
    
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Try extracting text directly
            page_text = page.get_text()
            
            # If text is too short, try OCR
            if len(page_text.strip()) < 50:
                # print(f"  Page {page_num+1} seems to be an image. Attempting OCR...")
                try:
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    image = Image.open(io.BytesIO(img_data))
                    ocr_text = pytesseract.image_to_string(image, lang='heb+eng')
                    page_text = ocr_text
                except Exception as e:
                    # print(f"  OCR failed: {e}")
                    pass
            
            text_content.append(page_text)
            
        full_text = "\n\n".join(text_content)
        return full_text
        
    except Exception as e:
        print(f"OCR Service Error: {e}")
        return ""
