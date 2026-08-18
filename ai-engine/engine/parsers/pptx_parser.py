"""
PowerPoint Parser — Extracts text and embedded images from PPTX files.

Uses python-pptx for structured XML parsing.
No OCR needed — PowerPoint stores text as structured data.
"""

import io
import logging
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import BaseParser, ParseResult

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):

    @staticmethod
    def supported_extensions() -> List[str]:
        return ['.pptx']

    def parse(self, file_bytes: bytes, filename: str) -> ParseResult:
        logger.info(f"PptxParser: Processing '{filename}'...")
        text_parts = []
        images = []

        try:
            prs = Presentation(io.BytesIO(file_bytes))

            for slide_num, slide in enumerate(prs.slides):
                slide_text = []

                for shape in slide.shapes:
                    # --- Text extraction ---
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # --- Table extraction ---
                    if shape.has_table:
                        table = shape.table
                        table_rows = []
                        for row in table.rows:
                            row_cells = [cell.text.strip() for cell in row.cells]
                            table_rows.append("| " + " | ".join(row_cells) + " |")
                        
                        if table_rows:
                            # Insert Markdown header separator after first row
                            header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                            table_rows.insert(1, header_sep)
                            slide_text.append("\n".join(table_rows))

                    # --- Image extraction ---
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_blob = shape.image.blob
                            if image_blob:
                                images.append(image_blob)
                                logger.debug(f"  Slide {slide_num + 1}: extracted image")
                        except Exception as e:
                            logger.warning(f"  Image extraction failed on slide {slide_num + 1}: {e}")

                text_parts.append(f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text))

        except Exception as e:
            logger.error(f"PptxParser Error for '{filename}': {e}")
            return ParseResult(text="", images=[], metadata={"error": str(e)})

        full_text = "\n\n".join(text_parts)
        logger.info(f"PptxParser: Extracted {len(full_text)} chars, {len(images)} images from '{filename}'")

        return ParseResult(
            text=full_text,
            images=images,
            metadata={"slide_count": len(text_parts), "image_count": len(images)}
        )

    def parse_in_chunks(self, file_bytes: bytes, filename: str, chunk_size: int = 1):
        logger.info(f"PptxParser: Processing '{filename}' in chunks of {chunk_size} slides...")
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            total_slides = len(prs.slides)
            
            for chunk_start in range(0, total_slides, chunk_size):
                text_parts = []
                images = []
                chunk_end = min(chunk_start + chunk_size, total_slides)
                
                for slide_num in range(chunk_start, chunk_end):
                    slide = prs.slides[slide_num]
                    slide_text = []
                    
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                        
                        if shape.has_table:
                            table = shape.table
                            table_rows = []
                            for row in table.rows:
                                row_cells = [cell.text.strip() for cell in row.cells]
                                table_rows.append("| " + " | ".join(row_cells) + " |")
                            if table_rows:
                                header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                                table_rows.insert(1, header_sep)
                                slide_text.append("\n".join(table_rows))
                        
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image_blob = shape.image.blob
                                if image_blob:
                                    images.append(image_blob)
                            except Exception as e:
                                logger.warning(f"  Image extraction failed on slide {slide_num + 1}: {e}")
                                
                    text_parts.append(f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text))
                
                yield ParseResult(
                    text="\n\n".join(text_parts),
                    images=images,
                    metadata={"chunk_start": chunk_start + 1, "chunk_end": chunk_end, "total_slides": total_slides}
                )
        except Exception as e:
            logger.error(f"PptxParser chunking Error for '{filename}': {e}")
            yield ParseResult(text="", images=[], metadata={"error": str(e)})
